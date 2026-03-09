#!/usr/bin/env python3
"""
pptx_to_accessible_html.py
Converts PowerPoint (.pptx) files to accessible HTML, preserving:
  - Heading structure (slide titles → <h1>, content titles → <h2>/<h3>)
  - All text content with proper semantics
  - Images with SME-authored alt text from PowerPoint
  - Slide speaker notes
  - Tables with proper <th>/<td> structure

Usage:
  Single file:  python pptx_to_accessible_html.py presentation.pptx
  Output name:  python pptx_to_accessible_html.py presentation.pptx -o output.html
  Batch folder: python pptx_to_accessible_html.py ./slides_folder/
  With notes:   python pptx_to_accessible_html.py presentation.pptx --include-notes

Requirements:
  pip install python-pptx
"""

import argparse
import base64
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx is not installed.")
    print("Install it with:  pip install python-pptx")
    sys.exit(1)


# ── Helpers ──────────────────────────────────────────────────────────────────

def get_alt_text(shape) -> str:
    """Extract SME-authored alt text from a shape's XML (the 'descr' attribute)."""
    try:
        # Alt text lives in <p:cNvPr descr="..."> inside nvPicPr or nvSpPr
        nvPr = shape._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr")
        # Use the standard DrawingML namespace path
        cNvPr = shape._element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
        )
        # Fallback: walk up through nvPicPr / nvSpPr
        for tag in [
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr",
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr",
            "{http://schemas.openxmlformats.org/presentationml/2006/main}nvGrpSpPr",
        ]:
            container = shape._element.find(f".//{tag}")
            if container is not None:
                cNvPr = container.find(
                    "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr"
                )
                if cNvPr is not None:
                    break

        # Most reliable: search all cNvPr elements under the shape
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for cNvPr in shape._element.iter(f"{{{ns}}}cNvPr"):
            descr = cNvPr.get("descr", "").strip()
            if descr:
                return descr

        # Also check the drawing namespace used for grouped/linked images
        ns2 = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for cNvPr in shape._element.iter(f"{{{ns2}}}cNvPr"):
            descr = cNvPr.get("descr", "").strip()
            if descr:
                return descr

        return ""
    except Exception:
        return ""


def get_alt_text_and_decorative(shape) -> tuple:
    """
    Returns (alt_text, is_decorative).

    PowerPoint 365 stores the decorative flag as an XML extension element:
      <adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2016/decorative" val="1"/>
    inside <a:extLst> under <p:cNvPr>. Older versions (and LibreOffice) may
    instead set decorative="1" directly as an attribute on <p:cNvPr>.
    We check both.

    Logic:
      - decorative flag found (either format) → is_decorative=True, alt=""
      - descr attribute is non-empty → is_decorative=False, alt=descr
      - neither → is_decorative=False, alt="" (caller falls back to generated text)
    """
    DECORATIVE_NS = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    try:
        # ── Method 1: PowerPoint 365 extLst extension ──
        for dec_el in shape._element.iter(f"{{{DECORATIVE_NS}}}decorative"):
            if dec_el.get("val", "0") == "1":
                return ("", True)

        # ── Method 2: direct attribute on cNvPr (older PowerPoint / LibreOffice) ──
        for el in shape._element.iter():
            if el.tag.endswith("}cNvPr") or el.tag == "cNvPr":
                if el.get("decorative", "0") == "1":
                    return ("", True)
                descr = el.get("descr", "").strip()
                if descr:
                    return (descr, False)

        return ("", False)
    except Exception:
        return ("", False)


def get_alt_text_reliable(shape) -> str:
    """
    More robust alt-text extractor that searches ALL cNvPr elements
    in the shape's XML tree regardless of namespace prefix used by the file.
    """
    try:
        # Iterate every element in the shape's XML subtree
        for el in shape._element.iter():
            # Match any cNvPr tag regardless of namespace
            if el.tag.endswith("}cNvPr") or el.tag == "cNvPr":
                descr = el.get("descr", "").strip()
                if descr:
                    return descr
        return ""
    except Exception:
        return ""


def image_to_data_uri(image_blob: bytes, content_type: str) -> str:
    """Convert raw image bytes to a base64 data URI for embedding in HTML."""
    b64 = base64.b64encode(image_blob).decode("ascii")
    return f"data:{content_type};base64,{b64}"


def escape_html(text: str) -> str:
    """Minimal HTML escaping."""
    return (
        text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
    )


MONOSPACE_FONTS = {
    "consolas", "courier", "courier new", "lucida console", "monaco",
    "menlo", "source code pro", "fira code", "fira mono",
    "dejavu sans mono", "liberation mono", "roboto mono",
    "jetbrains mono", "inconsolata", "cascadia code", "cascadia mono",
    "sf mono", "hack", "droid sans mono", "ubuntu mono",
    "noto sans mono", "ibm plex mono", "anonymous pro",
    "ocr a", "ocr b", "andale mono", "lucida sans typewriter",
}


def is_code_shape(shape) -> bool:
    """Detect code content via shape name convention or monospace font usage."""
    try:
        if shape.name and "code" in shape.name.lower():
            return True
    except Exception:
        pass
    if not shape.has_text_frame:
        return False
    mono_runs = 0
    total_runs = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            total_runs += 1
            font_name = run.font.name
            if font_name and font_name.strip().lower() in MONOSPACE_FONTS:
                mono_runs += 1
    return total_runs > 0 and mono_runs == total_runs


def render_code_block(tf) -> str:
    """Render a text frame as a <pre><code> block, preserving whitespace."""
    lines = []
    for para in tf.paragraphs:
        # Preserve original text including leading whitespace
        lines.append(escape_html(para.text))
    code_text = "\n".join(lines)
    return f"<pre><code>{code_text}</code></pre>"


def is_title_placeholder(shape) -> bool:
    """Return True if the shape is a title or centered-title placeholder."""
    from pptx.enum.text import PP_ALIGN
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        # Placeholder type 1 = TITLE, 3 = CENTER_TITLE, 15 = SUBTITLE
        return ph.type in (1, 3)
    except Exception:
        return False


def is_subtitle_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type == 15  # SUBTITLE only (not BODY=2, which is main content)
    except Exception:
        return False


def is_slide_number_placeholder(shape) -> bool:
    try:
        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type == 13  # SLIDE_NUMBER
    except Exception:
        return False


def is_running_header(shape, deck_title: str) -> bool:
    """Return True if shape is a body placeholder whose text matches the deck title."""
    if not deck_title:
        return False
    try:
        ph = shape.placeholder_format
        if ph is None or ph.type != 2:  # BODY
            return False
        if shape.has_text_frame:
            return shape.text_frame.text.strip() == deck_title
    except Exception:
        pass
    return False


def shape_has_text(shape) -> bool:
    return shape.has_text_frame and shape.text_frame.text.strip()


_CITATION_RE = re.compile(r'^[\d,\s\-–]+$')


def render_paragraph_runs(para, bracket_refs: bool = True) -> str:
    """Render a paragraph's runs to HTML, preserving superscript/subscript.

    If bracket_refs is True, superscript runs that look like citation numbers
    (digits, commas, spaces, hyphens) are wrapped in brackets: [1,2,3].
    """
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    if not para.runs:
        return escape_html(para.text.strip())
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        safe = escape_html(text)
        rPr = run._r.find(f"{A_NS}rPr")
        if rPr is not None:
            baseline = rPr.get("baseline")
            if baseline:
                val = int(baseline)
                if val > 0:
                    if bracket_refs and _CITATION_RE.match(text.strip()):
                        safe = f"<sup>[{safe}]</sup>"
                    else:
                        safe = f"<sup>{safe}</sup>"
                elif val < 0:
                    safe = f"<sub>{safe}</sub>"
        parts.append(safe)
    return "".join(parts).strip()


def _detect_list_type(para) -> str:
    """Detect paragraph list type: 'ol', 'ul', or '' (not a list item)."""
    A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        if para._p.pPr is not None:
            if para._p.pPr.find(f"{A_NS}buAutoNum") is not None:
                return "ol"
            if para._p.pPr.find(f"{A_NS}buChar") is not None:
                return "ul"
    except Exception:
        pass
    return ""


def render_text_frame(tf, base_heading_level: int = 3, bracket_refs: bool = True) -> str:
    """
    Render a text frame to HTML.
    Paragraphs that look like headings (bold, larger font) become <hN>.
    Everything else becomes <p>.
    """
    html_parts = []
    list_tag = ""  # "", "ul", or "ol"

    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            continue

        safe = render_paragraph_runs(para, bracket_refs=bracket_refs)

        # Detect bullet / numbered list paragraph
        para_list_type = _detect_list_type(para)

        # Detect heading-like paragraph (bold + larger font in first run)
        is_heading = False
        try:
            if para.runs:
                first_run = para.runs[0]
                font_size = first_run.font.size
                is_bold = first_run.font.bold
                if is_bold and font_size and font_size >= Pt(14):
                    is_heading = True
        except Exception:
            pass

        if is_heading:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            level = min(base_heading_level, 6)
            html_parts.append(f"<h{level}>{safe}</h{level}>")
        elif para_list_type:
            # Switch list type if changing between ul and ol
            if list_tag and list_tag != para_list_type:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            if not list_tag:
                html_parts.append(f"<{para_list_type}>")
                list_tag = para_list_type
            html_parts.append(f"  <li>{safe}</li>")
        else:
            if list_tag:
                html_parts.append(f"</{list_tag}>")
                list_tag = ""
            html_parts.append(f"<p>{safe}</p>")

    if list_tag:
        html_parts.append(f"</{list_tag}>")

    return "\n".join(html_parts)


def render_table(shape) -> str:
    """Render a table shape to an accessible HTML <table>."""
    table = shape.table
    rows_html = []
    for row_idx, row in enumerate(table.rows):
        cells_html = []
        for cell in row.cells:
            text = escape_html(cell.text.strip())
            tag = "th" if row_idx == 0 else "td"
            scope = ' scope="col"' if row_idx == 0 else ""
            cells_html.append(f"<{tag}{scope}>{text}</{tag}>")
        rows_html.append("<tr>" + "".join(cells_html) + "</tr>")

    caption = ""
    try:
        if shape.name:
            caption = f"<caption>{escape_html(shape.name)}</caption>"
    except Exception:
        pass

    return (
        '<table>\n'
        f'  {caption}\n'
        '  <thead>\n'
        f'    {rows_html[0]}\n'
        '  </thead>\n'
        '  <tbody>\n'
        + "\n".join(f"    {r}" for r in rows_html[1:])
        + "\n  </tbody>\n</table>"
    )


def merge_adjacent_lists(html: str) -> str:
    """Merge consecutive same-type list closings/openings with only whitespace between."""
    html = re.sub(r'</ul>\s*<ul>', '', html)
    html = re.sub(r'</ol>\s*<ol>', '', html)
    return html


# ── Core conversion ───────────────────────────────────────────────────────────

def get_slide_dimensions(prs):
    """Return slide canvas width and height in EMUs."""
    return prs.slide_width, prs.slide_height


def convert_slide(slide, slide_number: int, include_notes: bool, slide_width=None, slide_height=None, deck_title: str = "", bracket_refs: bool = True) -> tuple:
    """Convert a single slide, returning (title_text, content_html_parts).

    The caller is responsible for section wrapping and title deduplication.
    """
    parts = []
    title_text = ""
    content_shapes = []

    # Separate title from content; skip slide numbers and running headers
    for shape in slide.shapes:
        if is_title_placeholder(shape) and shape_has_text(shape):
            title_text = shape.text_frame.text.strip()
        elif is_slide_number_placeholder(shape):
            continue
        elif is_running_header(shape, deck_title):
            continue
        else:
            content_shapes.append(shape)

    # Process content shapes
    for shape in content_shapes:
        # ── Image ──
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            alt, is_decorative = get_alt_text_and_decorative(shape)
            if is_decorative:
                parts.append(f'  <!-- Decorative image skipped (slide {slide_number}) -->')
                continue
            if not alt:
                alt = f"Image on slide {slide_number}"
            try:
                img = shape.image
                data_uri = image_to_data_uri(img.blob, img.content_type)
                px_w = round(shape.width / 914400 * 96)
                px_h = round(shape.height / 914400 * 96)
                size_style = f' style="width:{px_w}px; height:{px_h}px; max-width:100%; height:auto;"'
                parts.append(
                    f'  <figure>\n'
                    f'    <img src="{data_uri}" alt="{escape_html(alt)}"{size_style}>\n'
                    f'  </figure>'
                )
            except Exception as e:
                parts.append(f'  <!-- Image could not be embedded: {e} -->')

        # ── Table ──
        elif shape.has_table:
            parts.append(render_table(shape))

        # ── Code block (monospace font or shape named "code") ──
        elif shape.has_text_frame and is_code_shape(shape):
            rendered = render_code_block(shape.text_frame)
            if rendered:
                parts.append(rendered)

        # ── Text frame ──
        elif shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if is_subtitle_placeholder(shape):
                parts.append(f"  <p class=\"subtitle\">{escape_html(text)}</p>")  # type 15 only
            else:
                rendered = render_text_frame(shape.text_frame, base_heading_level=3, bracket_refs=bracket_refs)
                if rendered:
                    parts.append(rendered)

        # ── Group shape — recurse for nested images ──
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    alt, is_decorative = get_alt_text_and_decorative(child)
                    if is_decorative:
                        continue
                    if not alt:
                        alt = f"Image on slide {slide_number}"
                    try:
                        img = child.image
                        data_uri = image_to_data_uri(img.blob, img.content_type)
                        size_style = ""
                        if slide_width and slide_height and slide_width > 0 and slide_height > 0:
                            px_w = round(child.width / 914400 * 96)
                            px_h = round(child.height / 914400 * 96)
                            size_style = f' style="width:{px_w}px; height:{px_h}px; max-width:100%; height:auto;"'
                        parts.append(
                            f'  <figure>\n'
                            f'    <img src="{data_uri}" alt="{escape_html(alt)}"{size_style}>\n'
                            f'  </figure>'
                        )
                    except Exception:
                        pass

    # ── Speaker notes ──
    if include_notes:
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(
                    f'  <aside class="speaker-notes">\n'
                    f'    <h3>Speaker Notes</h3>\n'
                    f'    <p>{escape_html(notes_text)}</p>\n'
                    f'  </aside>'
                )
        except Exception:
            pass

    return (title_text, parts)


def convert_pptx(input_path: Path, output_path: Path, include_notes: bool, bracket_refs: bool = True) -> None:
    """Convert a .pptx file to an accessible HTML file."""
    prs = Presentation(str(input_path))
    title = input_path.stem.replace("_", " ").replace("-", " ").title()

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Extract deck title from slide 1 to filter running headers on other slides
    deck_title = ""
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if is_title_placeholder(shape) and shape_has_text(shape):
                deck_title = shape.text_frame.text.strip()
                break

    # Collect (title, content_parts, slide_number) for each slide
    slide_data = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_title, parts = convert_slide(slide, i, include_notes, slide_width, slide_height, deck_title=deck_title, bracket_refs=bracket_refs)
        slide_data.append((slide_title, parts, i))

    # Group consecutive slides that share the same title
    # Each group becomes one <section> with one <h2>
    groups = []  # list of (title, first_slide_number, combined_parts)
    for slide_title, parts, slide_num in slide_data:
        effective_title = slide_title if slide_title else f"Slide {slide_num}"
        if groups and groups[-1][0] == effective_title:
            # Same title as previous — append content to existing group
            groups[-1][2].extend(parts)
        else:
            groups.append((effective_title, slide_num, list(parts)))

    # Build section HTML for each group, merging adjacent lists
    slide_sections = []
    for group_title, first_slide_num, content_parts in groups:
        section_lines = [f'<section aria-label="Slide {first_slide_num}">']
        section_lines.append(f"  <h2>{escape_html(group_title)}</h2>")
        content_html = "\n".join(content_parts)
        content_html = merge_adjacent_lists(content_html)
        section_lines.append(content_html)
        section_lines.append("</section>")
        slide_sections.append("\n".join(section_lines))

    slides_html = "\n\n".join(slide_sections)
    total = len(prs.slides)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{escape_html(title)}</title>
  <style>
    /* ── Base ── */
    *, *::before, *::after {{ box-sizing: border-box; }}
    body {{
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
      font-size: 1rem;
      line-height: 1.6;
      color: #1a1a1a;
      background: #f5f5f5;
      margin: 0;
      padding: 2rem 1rem;
    }}
    /* ── Skip link ── */
    .skip-link {{
      position: absolute;
      top: -40px;
      left: 0;
      background: #005fcc;
      color: #fff;
      padding: 0.5rem 1rem;
      border-radius: 0 0 4px 0;
      text-decoration: none;
      font-weight: bold;
      z-index: 100;
    }}
    .skip-link:focus {{ top: 0; }}
    /* ── Layout ── */
    main {{ max-width: 900px; margin: 0 auto; }}
    h1 {{
      font-size: 2rem;
      margin-bottom: 0.25rem;
      color: #111;
    }}
    .deck-meta {{
      color: #555;
      margin-bottom: 2rem;
      font-size: 0.95rem;
    }}
    /* ── Slide sections ── */
    section {{
      background: #fff;
      border: 1px solid #ddd;
      border-radius: 8px;
      padding: 2rem;
      margin-bottom: 2rem;
      box-shadow: 0 1px 4px rgba(0,0,0,.06);
    }}
    section h2 {{
      font-size: 1.5rem;
      margin-top: 0;
      color: #003366;
      border-bottom: 2px solid #e0e0e0;
      padding-bottom: 0.5rem;
    }}
    section h3 {{
      font-size: 1.15rem;
      color: #1a1a1a;
    }}
    p {{ margin: 0.5rem 0; }}
    ul {{ padding-left: 1.5rem; margin: 0.5rem 0; }}
    li {{ margin: 0.25rem 0; }}
    /* ── Code blocks ── */
    pre {{
      background: #1e1e1e;
      color: #d4d4d4;
      border-radius: 6px;
      padding: 1rem 1.25rem;
      overflow-x: auto;
      margin: 1rem 0;
      font-size: 0.9rem;
      line-height: 1.5;
    }}
    pre code {{
      font-family: Consolas, "Courier New", "Fira Mono", monospace;
      white-space: pre;
    }}
    /* ── Images ── */
    figure {{
      margin: 1rem 0;
      text-align: center;
    }}
    figure img {{
      max-width: 100%;
      height: auto;
      border-radius: 4px;
      border: 1px solid #e0e0e0;
    }}
    /* ── Tables ── */
    table {{
      border-collapse: collapse;
      width: 100%;
      margin: 1rem 0;
      font-size: 0.95rem;
    }}
    caption {{
      font-weight: bold;
      margin-bottom: 0.5rem;
      text-align: left;
    }}
    th, td {{
      border: 1px solid #ccc;
      padding: 0.5rem 0.75rem;
      text-align: left;
    }}
    th {{
      background: #003366;
      color: #fff;
    }}
    tr:nth-child(even) td {{ background: #f9f9f9; }}
    /* ── Speaker notes ── */
    .speaker-notes {{
      margin-top: 1.5rem;
      padding: 1rem;
      background: #fffbea;
      border-left: 4px solid #f0c040;
      border-radius: 4px;
    }}
    .speaker-notes h3 {{
      margin-top: 0;
      font-size: 0.9rem;
      text-transform: uppercase;
      letter-spacing: 0.05em;
      color: #7a6000;
    }}
    /* ── Subtitle ── */
    p.subtitle {{
      font-size: 1.1rem;
      color: #444;
      font-style: italic;
    }}
  </style>
</head>
<body>
  <a href="#main-content" class="skip-link">Skip to main content</a>
  <main id="main-content">
    <h1>{escape_html(title)}</h1>
    <p class="deck-meta">Converted from <strong>{escape_html(input_path.name)}</strong> — {total} slide{"s" if total != 1 else ""}</p>

{slides_html}

  </main>
</body>
</html>
"""

    output_path.write_text(html, encoding="utf-8")
    print(f"OK  {input_path.name}  ->  {output_path}")


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint (.pptx) to accessible HTML, preserving alt text."
    )
    parser.add_argument(
        "input",
        help="Path to a .pptx file, or a folder containing .pptx files."
    )
    parser.add_argument(
        "-o", "--output",
        help="Output HTML file path (single-file mode only).",
        default=None
    )
    parser.add_argument(
        "--include-notes",
        action="store_true",
        help="Include speaker notes in the HTML output."
    )
    parser.add_argument(
        "--no-bracket-refs",
        action="store_true",
        help="Disable bracketing of superscript citation numbers (e.g. [1,2,3])."
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    bracket_refs = not args.no_bracket_refs

    # ── Batch mode: folder ──
    if input_path.is_dir():
        pptx_files = sorted(input_path.glob("*.pptx"))
        if not pptx_files:
            print(f"No .pptx files found in: {input_path}")
            sys.exit(1)
        for pptx_file in pptx_files:
            out = pptx_file.with_suffix(".html")
            convert_pptx(pptx_file, out, args.include_notes, bracket_refs=bracket_refs)
        print(f"\nDone. {len(pptx_files)} file(s) converted.")
        return

    # ── Single-file mode ──
    if not input_path.exists():
        print(f"ERROR: File not found: {input_path}")
        sys.exit(1)
    if input_path.suffix.lower() != ".pptx":
        print(f"ERROR: Expected a .pptx file, got: {input_path.suffix}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_suffix(".html")

    convert_pptx(input_path, output_path, args.include_notes, bracket_refs=bracket_refs)
    print("Done.")


if __name__ == "__main__":
    main()
